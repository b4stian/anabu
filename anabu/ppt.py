#!/usr/bin/env python3

# ------------------------------------------------
# imports

try:
    import charts
    import density
    import interface
    import photo
    import pinholes
except:
    from anabu import charts
    from anabu import interface
    from anabu import photo
    from anabu import density
    from anabu import pinholes

import os

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT

# ------------------------------------------------
# variables

PATH_MASTER_PPT = r"ppt/PPT_master.pptx"
# This module will fail if wrong PPT_master is used.

# ------------------------------------------------
# functions/classes


class Pptx:
    def __init__(self, master_pptx: str) -> None:
        try:
            self.presentation = Presentation(master_pptx)
        except:
            interface.logging.info("Master PPTX not found. Using file dialog.")
            master_pptx = interface.filedialog.askopenfilename(
                filetypes=[("PPTX Files", ".pptx")],
                title="Select master PPTX file.",
            )
            interface.logging.info(f"Selected photo file using dialog: {master_pptx}")
            self.presentation = Presentation(master_pptx)
        interface.results.add_result(
            variable="PPTX_master",
            parameter="master PPTX file",
            value=master_pptx,
        )
        interface.logging.info(
            f"Initiated PPTX with results from master {master_pptx}."
        )

    def generate_title_slide(self) -> None:
        self.title_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[0]
        )
        self.title_slide.placeholders[0].text = interface.results.sample_name["value"]
        self.title_slide.placeholders[1].text = (
            "anabu version " + interface.results.version["value"]
        )
        self.title_slide.placeholders[
            13
        ].text = "Evaluation results for light table photo"
        interface.logging.info("Generated title slide of PPTX.")

    def generate_cropped_photo_slide(self) -> None:
        self.cropped_photo_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[10]
        )
        self.cropped_photo_slide.placeholders[
            0
        ].text = f"Cropped photo of {interface.results.sample_name['value']}"
        self.cropped_photo_slide.placeholders[13].insert_picture(
            os.path.splitext(photo.photo.photo_path)[0] + "_scale_axes.png"
        )
        interface.logging.info("Generated slide with cropped photo.")

    def generate_distribution_plot_slide(self) -> None:
        self.distribution_plot_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[10]
        )
        self.distribution_plot_slide.placeholders[
            0
        ].text = f"Brightness distribution of {interface.results.sample_name['value']}"
        self.distribution_plot_slide.placeholders[13].insert_picture(
            os.path.splitext(photo.photo.photo_path)[0] + "_distribution_plot.png"
        )
        interface.logging.info("Generated slide with distribution plots.")

    def generate_maskview_slide(self) -> None:
        if interface.user_settings.maskview["value"] == None:
            interface.logging.info("Did not generated slide with Maskview because it is switched off.")
            return
        self.maskview_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[11]
        )
        self.maskview_slide.placeholders[
            0
        ].text = f"Maskview of {interface.results.sample_name['value']}"
        self.maskview_slide.placeholders[13].insert_picture(
            os.path.splitext(photo.photo.photo_path)[0] + "_maskview.tiff"
        )
        interface.logging.info("Generated slide with Maskview.")

    def generate_results_slide(self) -> None:
        self.results_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[12]
        )
        self.results_slide.placeholders[
            0
        ].text = f"Results for {interface.results.sample_name['value']}"
        shape = self.results_slide.placeholders[14].insert_table(rows=12, cols=2)
        table = shape.table
        tbl = shape._element.graphic.graphicData.tbl
        style_id = "{8EC20E35-A176-4012-BC5E-935CFFF8708E}"
        tbl[0][-1].text = style_id
        cell = table.cell(0, 0)
        cell.text = "parameter"
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        cell = table.cell(0, 1)
        cell.text = "value"
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        cell = table.cell(1, 0)
        cell.text = interface.results.creation_date["parameter"]
        cell = table.cell(1, 1)
        cell.text = interface.results.creation_date["value"]
        cell = table.cell(2, 0)
        cell.text = "camera, F number, exposure time, ISO"
        cell = table.cell(2, 1)
        cell.text = f'{str(interface.results.Make["value"])} {str(interface.results.Model["value"])}, {str(interface.results.FNumber["value"])}, {str(interface.results.ExposureTime["value"])} s, {str(interface.results.ISOSpeedRatings["value"])}'
        cell = table.cell(3, 0)
        cell.text = "operator"
        cell = table.cell(3, 1)
        cell.text = interface.results.operator["value"]
        cell = table.cell(4, 0)
        cell.text = interface.results.number_pixels["parameter"]
        cell = table.cell(4, 1)
        cell.text = (
            str(interface.results.number_pixels["value"])
            + f" ({round(interface.results.percentage_nonmasked['value'],1)}% of total)"
        )
        cell = table.cell(5, 0)
        cell.text = "minimum, maximum brightness value"
        cell = table.cell(5, 1)
        cell.text = f'{str(interface.results.brightness_min["value"])}, {str(interface.results.brightness_max["value"])}'
        cell = table.cell(6, 0)
        cell.text = interface.results.brightness_peak["parameter"]
        cell = table.cell(6, 1)
        cell.text = str(interface.results.brightness_peak["value"])
        cell = table.cell(7, 0)
        cell.text = "brightness skew and kurtosis"
        cell = table.cell(7, 1)
        cell.text = f'{str(round(interface.results.brightness_skew["value"],3))}, {str(round(interface.results.brightness_kurtosis["value"],3))}'
        cell = table.cell(8, 0)
        cell.text = interface.results.brightness_mean["parameter"]
        cell = table.cell(8, 1)
        cell.text = str(round(interface.results.brightness_mean["value"], 3))
        cell = table.cell(9, 0)
        cell.text = "25th, 50th, 75th, 95th, 99th percentile"
        cell = table.cell(9, 1)
        cell.text = f'{str(interface.results.brightness_25["value"])}, {str(interface.results.brightness_median["value"])}, {str(interface.results.brightness_75["value"])}, {str(interface.results.brightness_95["value"])}, {str(interface.results.brightness_99["value"])}'
        cell = table.cell(10, 0)
        cell.text = "standard deviation and dispersity"
        cell = table.cell(10, 1)
        cell.text = f'{str(round(interface.results.brightness_stdev["value"],3))}, {str(round(interface.results.brightness_dispersity["value"],3))}'
        cell = table.cell(11, 0)
        cell.text = "optical density (calibration)"
        cell = table.cell(11, 1)
        cell.text = f'{str(interface.results.optical_density["value"])} ({str(interface.results.calibration_density_name["value"])})'
        interface.logging.info("Generated slide with results table.")

    def generate_pinholes_slide(self) -> None:
        if not interface.user_settings.pinholes["value"]:
            return
        im = Image.open(os.path.splitext(photo.photo.photo_path)[0] + "_pinholes.png")
        x, y = im.size
        new_x = int(y * 335 / 150)
        new_im = Image.new("RGBA", (new_x, y), (0, 0, 0, 0))
        new_im.paste(im, (int((new_x - x) / 2), 0))
        new_im.save(
            os.path.splitext(photo.photo.photo_path)[0] + "_pinholes_resized.png"
        )
        self.pinholes_slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[10]
        )
        self.pinholes_slide.placeholders[
            0
        ].text = f"Pinhole analysis of {interface.results.sample_name['value']} ({interface.results.calibration_pinholes_name['value']})"
        placeholder = self.pinholes_slide.placeholders[13]
        placeholder = placeholder.insert_picture(
            os.path.splitext(photo.photo.photo_path)[0] + "_pinholes_resized.png"
        )
        interface.logging.info("Generated slide with marked pinholes.")

    def save_presentation(self) -> None:
        self.presentation.save(
            os.path.splitext(photo.photo.photo_path)[0] + "_results.pptx"
        )
        interface.logging.info(
            f"Saved PPTX with results to {os.path.splitext(photo.photo.photo_path)[0] + '_results.pptx'}."
        )


def run_pptx() -> None:
    if not interface.user_settings.create_ppt["value"]:
        return
    prs = Pptx(PATH_MASTER_PPT)
    prs.generate_title_slide()
    interface.Gui.update_progress_bar()
    prs.generate_cropped_photo_slide()
    interface.Gui.update_progress_bar()
    prs.generate_distribution_plot_slide()
    interface.Gui.update_progress_bar()
    prs.generate_results_slide()
    interface.Gui.update_progress_bar()
    prs.generate_pinholes_slide()
    interface.Gui.update_progress_bar()
    prs.generate_maskview_slide()
    interface.Gui.update_progress_bar()
    prs.save_presentation()
    interface.Gui.update_progress_bar()


# ------------------------------------------------
# executions

is_main = __name__ == "__main__"

if is_main:
    interface.run_interface()
    photo.run_photo()
    density.run_density()
    pinholes.run_pinholes()
    charts.run_charts()
    run_pptx()


# for shape in self.title_slide.placeholders:
#             print('%d %s' % (shape.placeholder_format.idx, shape.name))
